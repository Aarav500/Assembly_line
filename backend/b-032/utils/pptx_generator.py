from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


def _add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    # Style
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    if subtitle:
        slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(20)
    return slide


def _add_content_slide(prs, title, body):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    paragraphs = [p.strip() for p in body.replace('\r', '').split('\n') if p.strip()]
    if not paragraphs:
        paragraphs = [body.strip()] if body.strip() else [" "]

    # First paragraph as bullet
    p = tf.paragraphs[0]
    p.text = paragraphs[0]
    p.level = 0
    p.font.size = Pt(18)

    for para in paragraphs[1:]:
        pp = tf.add_paragraph()
        pp.text = para
        pp.level = 0
        pp.font.size = Pt(18)

    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    return slide


def _brand_colors(prs):
    # Not strictly necessary; keep default template.
    pass


def generate_pitch_deck_pptx(data, slides):
    prs = Presentation()

    company = (data.get('company_name') or 'Your Company').strip() or 'Your Company'
    tagline = (data.get('tagline') or '').strip()

    _add_title_slide(prs, company, tagline)

    for s in slides[1:]:  # first slide already created as title
        title = s.get('title', '') or 'Slide'
        content = s.get('content', '') or ''
        _add_content_slide(prs, title, content)

    # Footer on all slides
    for slide in prs.slides:
        left = Inches(0.5)
        top = Inches(6.9)
        width = Inches(9)
        height = Inches(0.3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = (company + (f" — {tagline}" if tagline else ''))
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

